# app/domains/knowledge/service.py
"""Internal document ingestion helpers.

This module is kept small because the current router does not expose document
upload endpoints yet. The functions below are safe building blocks for the
future knowledge-base API and keep the merged code importable.
"""

from __future__ import annotations

import io, re, os, subprocess, tempfile
from datetime import datetime
from pathlib import Path
from typing import Any, Optional
from langchain_text_splitters import RecursiveCharacterTextSplitter

from app.domains.knowledge.agent_utils import chroma_client, get_collection
from app.utils.time_utils import now_kst


_splitter = RecursiveCharacterTextSplitter(
    chunk_size=800,
    chunk_overlap=100,
    separators=["\n\n", "\n", "。", ". ", " ", ""],
)

def _collection_name(workspace_id: int) -> str:
    """Return the workspace-scoped Chroma collection name."""
    return f"ws_{workspace_id}_docs"

def _extract_pdf(file_bytes: bytes) -> str:
    """Extract plain text from a PDF file."""
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(file_bytes))
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text and text.strip():
            pages.append(text.strip())
    return "\n\n".join(pages)

def _extract_pptx(file_bytes: bytes) -> str:
    """
    PPT/PPTX -> 텍스트.

    슬라이드별 도형(shape) 텍스트 프레임 순회.
    "[슬라이드 N]" 헤더 삽입 -> 검색 결과에서 몇 번째 슬라이드인지 참조 가능.
    차트/이미지 속 텍스트는 추출 불가 -> vision 도메인 OCR 필요.
    """
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        if texts:
            slides.append(f"[슬라이드 (i)]\n" + "\n".join(texts))
    return "\n\n".join(slides)

def _extract_html(file_bytes: bytes) -> str:
    """
    HTML -> 텍스트.

    제거 태그: script/style(코드 노이즈), nav/header/footer(반복 메뉴 텍스트).
    연속 줄바꿈 3개 이상 -> 2개로 압축해 빈 청크 방지.
    """
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(file_bytes, "html.parser")
    for tag in soup(["script", "style", "nav", "header", "footer"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    return re.sub(fr"\n{3,}", "\n\n", text).strip()

def _extract_md(file_bytes: bytes) -> str:
    """
    Markdown -> 텍스트.

    별도 라이브러리 불필요 - UTF-8 텍스트 그대로 반환.
    마크다운 문법 기호(#, *, `, -)는 의미 있는 컨텍스트이므로 제거하지 않음.
    LLM이 마크다운을 이해하므로 오히려 구조 보존이 검색 품질에 유리.
    """
    return file_bytes.decode('utf-8', errors='replace')

def _extract_docx(file_bytes: bytes) -> str:
    """
    DOCX -> 텍스트.

    단락(paragraph) + 표(table) 두 가지 요소 순회.
    표는 행 단위로 셀을 탭으로 구분해 검색 가능한 텍스트로 변환.
    이미지/도형 안 텍스트는 추출 불가 -> vision 도메인 OCR 필요.

    주의: 구형 .doc 포맷은 python-docx 미지원 -> 422 반환
    """
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    parts = []

    # 단락 추출
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    # 표 추출 - 행 단위로 셀을 탭으로 이어붙임
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)

    return "\n\n".join(parts)

def _extrac_doc_legacy(file_bytes: bytes) -> str:
    """
    구형 .doc -> LibreOffice로 .docx 변환 후 python-docx로 추출.

    python-docx는 구현 .doc 포맷 미지원.
    LibreOffice --headless 모드로 .docx 변환 후 _extract_docx()에 위임.
    vision 도메인의 PPTX->PDF 변환과 동일한 LibreOffice 의존성.
    변환 타임아웃 30초 - 대용량 파일은 초과 가능.
    """
    with tempFile.TemporaryDirectory() as tmpdir:
        doc_path = os.path.join(tmdir, "input, doc")
        with open(doc_path, "wb") as f:
            f.write(file_bytes)

        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "docx", "--outdir", tmpdir, doc_path],
            capture_output=True,
            timeout=30,
        )
        if result.returncode != 0:
            raise ValueError("구형 .doc 변환 실패. 서버에 LibreOffice가 설치되어 있는지 확인하세요.")

        docx_path = os.path.join(tmpdir, "input-docx")
        with open(docx_path, "rb") as f:
            return _extract_docx(f.read())

def _extract_xlsx(file_bytes: bytes) -> str:
    """
    XLSX -> 텍스트.

    시트별로 "[시트명]" 헤더를 붙여 구분.
    수식 캐싱값 없는 파일은 LibreOffice로 재계산 후 재시도.
    LibreOffice 미설치 시 수식 문자열 fallback.
    """
    from openpyxl import load_workbook

    def _read(data: bytes) -> str:
        wb = load_workbook(io.BytesIO(data), data_only=True)
        sheets = []
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    rows.append("\t".join(cells))
            if rows:
                sheets.append(f"[{sheet.title}]\n" + "\n".join(rows))
        return "\n\n".join(sheets)

    text = _read(file_bytes)

    # None 셀 비율이 높으면 LibreOffice로 재계산 시도
    non_ratio = text.count("\t\t") / max(text.count("\t"), 1)
    if non_ratio > 0.3:
        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                xlsx_path = os.path.join(tmpdir, "input.xlsx")
                with open(xlsx_path, "wb") as f:
                    f.write(file_bytes)
                subprocess.run(
                    ["libreoffice", "--headless", "--convert-to", "xlsx", "--outdir", tmpdir, xlsx_path],
                    capture_output=True, timeout=30
                )
                out_path = os.path.join(tmpdir, "input.xlsx")
                with open(out_path, "rb") as f:
                    text = _read(f.read())
        except Exception:
            pass # LibreOffice 없으면 원래 결과 그대로 사용

    return text

def _extract_xls_legacy(file_bytes: bytes) -> str:
    """
    구형 .xls -> LibreOffice로 .xlsx 변환 후 _extract_xlsx()에 위임

    xlrd 대신 LibreOffice를 쓰는 이유:
        1. xlrd는 수식 캐싱값을 못 읽는 케이스 동일하게 존재.
        2. LibreOffice 변환 시 수식 재계산까지 처리.
        3. xlrd 의존성 제거.
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        xls_path = os.path.join(tmpdir, "input.xls")
        with open(xls_path, "wb") as f:
            f.write(file_bytes)

        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "xlsx", "--outdir", tmpdir, xls_path],
            capture_output=True,
            timeout=30,
        )
        if result.returncode != 0:
            raise ValueError("구형 .xls 변환 실패. 서버에 LibreOffice가 설치되어 있는지 확인하세요.")

        xlsx_path = os.path.join(tmpdir, "input.xlsx")
        with open(xlsx_path, "rb") as f:
            return _extract_xlsx(f.read())

def ingest_document(
    workspace_id: int,
    filename: str,
    file_bytes: bytes,
    file_type: str,
    title: Optional[str] = None,
) -> dict:
    """
    문서 수집 전체 파이프라인.

    Args:
        workspace_id: 워크스페이스 ID, 컬렉션 격리 키.
        filename: 웝본 파일명. doc_id 및 메타데이터에 사용.
        file_bytes: 업로드된 파일 바이너리.
        file_type: "pdf" | "pptx" | "html".
        title: 문서 제목. None 이면 filename 사용.

    Returns:
        {"doc_id": str, "chunks": int, "title": str}

    중복 업로드 처리:
        doc_id = "{workspace_id}_{filename}" 고정
        청크 ID도 "{doc_id}_{chunk_id}" 고정.
        -> 같은 파일 재업로드 시 upsert가 덮어씀. 벡터 중복 없음
    """
    # 1단계: 텍스트 추출
    if file_type == "pdf":
        raw_text = _extract_pdf(file_bytes)
    elif file_type == "pptx":
        raw_text = _extract_pptx(file_bytes)
    elif file_type == "html":
        raw_text = _extract_html(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {file_type} pdf | pptx | html 만 가능")

    if not raw_text.strip():
        raise ValueError(
            "텍스트를 추출할 수 없습니다. "
            "스캔 이미지 PDF는 vision 도메인 OCR(/api/v1/vision)을 사용하세요."
        )

    # 2단계: 청크 분할
    # 결과: 각 청크 <= 800자, 인접 청크 간 100자 overlap
    chunks = _splitter.split_text(raw_text)

    # 3단계: 메타데이터 구성
    doc_id = f"{workspace_id}_{filename}"
    title = title or filename
    uploaded_at = now_kst().isoformat()

    ids = [f"{doc_id}_chunk{i}" for i in range(len(chunks))]
    metadatas = [
        {
            "workspace_id": workspace_id,
            "doc_id": doc_id,
            "title": title,
            "filename": filename,
            "file_type": file_type,
            "chunk_index": i,
            "total_chunks": len(chunks),
            "uploaded_at": uploaded_at,
        }
        for i in range(len(chunks))
    ]

    # 4단계: ChromaDB upsert
    # add() 대신 upsert() 이유:
    #   add()는 동일 ID 존재 시 에러 → 재업로드 불가
    #   upsert()는 있으면 덮어쓰고 없으면 삽입 → 재업로드 안전
    collection = get_collection(workspace_id)
    collection.upsert(documents=chunks, ids=ids, metadatas=metadatas)

    return {"doc_id": doc_id, "chunks": len(chunks), "title": title}
